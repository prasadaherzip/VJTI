# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours ──────────────────────────────────────────────────────
WHITE   = RGBColor(0xFF,0xFF,0xFF)
BLACK   = RGBColor(0x1A,0x1A,0x2E)
GREY    = RGBColor(0x55,0x55,0x66)
LGREY   = RGBColor(0xF2,0xF2,0xF5)
MGREY   = RGBColor(0xCC,0xCC,0xD8)
BLUE    = RGBColor(0x37,0x8A,0xDD)
LBLUE   = RGBColor(0xE6,0xF1,0xFB)
DBLUE   = RGBColor(0x0C,0x44,0x7C)
TEAL    = RGBColor(0x1D,0x9E,0x75)
LTEAL   = RGBColor(0xE1,0xF5,0xEE)
DTEAL   = RGBColor(0x08,0x50,0x41)
PURPLE  = RGBColor(0x53,0x4A,0xB7)
LPURPLE = RGBColor(0xEE,0xED,0xFE)
DPURPLE = RGBColor(0x3C,0x34,0x89)
RED     = RGBColor(0xE2,0x4B,0x4A)
LRED    = RGBColor(0xFC,0xEB,0xEB)
DRED    = RGBColor(0x79,0x1F,0x1F)
AMBER   = RGBColor(0xEF,0x9F,0x27)
LAMBER  = RGBColor(0xFA,0xEE,0xDA)
DAMBER  = RGBColor(0x63,0x38,0x06)

W = Inches(13.33)
H = Inches(7.5)
FONT = "Calibri"

# ── Helpers ───────────────────────────────────────────────────────
def new_prs():
    p = Presentation()
    p.slide_width  = W
    p.slide_height = H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, top, w, h, fill_color=None, line_color=None, line_pt=0.75):
    shp = slide.shapes.add_shape(1, l, top, w, h)
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_pt)
    else:
        shp.line.fill.background()
    return shp

def add_oval(slide, l, top, w, h, fill_color):
    shp = slide.shapes.add_shape(9, l, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp

# renamed from t() to ln() to avoid clash with loop variable t
def ln(text, size, bold=False, italic=False, color=BLACK,
       align=PP_ALIGN.LEFT, space_before=0):
    return dict(text=text, size=size, bold=bold, italic=italic,
                color=color, align=align, space_before=space_before)

def textbox(slide, l, top, w, h, lines, wrap=True):
    txb = slide.shapes.add_textbox(l, top, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(line.get("space_before", 0))
        r = p.add_run()
        r.text = line["text"]
        r.font.name = FONT
        r.font.size = Pt(line["size"])
        r.font.bold = line.get("bold", False)
        r.font.italic = line.get("italic", False)
        r.font.color.rgb = line.get("color", BLACK)
    return txb

# ── Reusable components ───────────────────────────────────────────
def top_bar(slide, color):
    add_rect(slide, 0, 0, W, Inches(0.22), fill_color=color)

def slide_tag(slide, label, fill, fg):
    shp = add_rect(slide, Inches(0.55), Inches(0.35),
                   Inches(3.4), Inches(0.32), fill_color=fill)
    shp.line.fill.background()
    tf = shp.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = FONT
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = fg

def slide_title(slide, text):
    textbox(slide, Inches(0.55), Inches(0.75), Inches(12.2), Inches(0.65),
            [ln(text, 26, bold=True, color=BLACK)])

def divider(slide, top_pos):
    add_rect(slide, Inches(0.55), top_pos, Inches(12.2), Pt(0.8), fill_color=MGREY)

def slide_num(slide, n, total):
    textbox(slide, Inches(12.2), Inches(0.28), Inches(0.9), Inches(0.28),
            [ln(f"{n} / {total}", 10, color=GREY, align=PP_ALIGN.RIGHT)])

def card_shape(slide, l, top, w, h, fill=LGREY):
    shp = add_rect(slide, l, top, w, h, fill_color=fill)
    return shp

def card_content(slide, l, top, w, h, heading, body_lines, cite=""):
    card_shape(slide, l, top, w, h)
    textbox(slide, l+Inches(0.15), top+Inches(0.12),
            w-Inches(0.3), Inches(0.38),
            [ln(heading, 13, bold=True, color=BLACK)])
    body_text = "\n".join(f"• {b}" for b in body_lines) if isinstance(body_lines, list) else body_lines
    textbox(slide, l+Inches(0.15), top+Inches(0.5),
            w-Inches(0.3), h-Inches(0.85),
            [ln(body_text, 11, color=GREY)])
    if cite:
        textbox(slide, l+Inches(0.15), top+h-Inches(0.28),
                w-Inches(0.3), Inches(0.26),
                [ln(cite, 8, italic=True, color=GREY)])

def stat_card(slide, l, top, w, h, stat_text, heading, body, cite=""):
    card_shape(slide, l, top, w, h)
    textbox(slide, l+Inches(0.15), top+Inches(0.1),
            w-Inches(0.3), Inches(0.5),
            [ln(stat_text, 28, bold=True, color=BLACK)])
    textbox(slide, l+Inches(0.15), top+Inches(0.58),
            w-Inches(0.3), Inches(0.35),
            [ln(heading, 12, bold=True, color=BLACK)])
    textbox(slide, l+Inches(0.15), top+Inches(0.92),
            w-Inches(0.3), h-Inches(1.1),
            [ln(body, 11, color=GREY)])
    if cite:
        textbox(slide, l+Inches(0.15), top+h-Inches(0.28),
                w-Inches(0.3), Inches(0.26),
                [ln(cite, 8, italic=True, color=GREY)])

def half_card(slide, l, top, w, h, heading, points, cite):
    card_shape(slide, l, top, w, h)
    textbox(slide, l+Inches(0.18), top+Inches(0.12),
            w-Inches(0.36), Inches(0.38),
            [ln(heading, 14, bold=True, color=BLACK)])
    divider(slide, top+Inches(0.58))
    body = "\n".join(f"• {p}" for p in points)
    textbox(slide, l+Inches(0.18), top+Inches(0.68),
            w-Inches(0.36), h-Inches(1.12),
            [ln(body, 12, color=GREY)])
    textbox(slide, l+Inches(0.18), top+h-Inches(0.32),
            w-Inches(0.36), Inches(0.28),
            [ln(cite, 8, italic=True, color=GREY)])

# ════════════════════════════════════════════════════════════════
prs = new_prs()

# ── SLIDE 1 — Cover ──────────────────────────────────────────────
s = blank(prs)
set_bg(s, WHITE)
top_bar(s, BLUE)
slide_num(s, 1, 8)

shp = add_rect(s, Inches(0.55), Inches(0.38), Inches(5.5), Inches(0.3), fill_color=LBLUE)
shp.line.fill.background()
tf = shp.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "IndiaNext Hackathon 2026  ·  K.E.S. Shroff College, Mumbai"
r.font.name = FONT; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = DBLUE

textbox(s, Inches(0.55), Inches(0.88), Inches(9), Inches(1.3),
        [ln("Third Eye", 52, bold=True, color=BLACK)])
textbox(s, Inches(0.55), Inches(2.1), Inches(10), Inches(0.5),
        [ln("Smart AI-Powered Cyber Defense Platform", 18, color=GREY)])
divider(s, Inches(2.72))
textbox(s, Inches(0.55), Inches(2.88), Inches(4), Inches(0.32),
        [ln("Team: Runtime Errors", 11, italic=True, color=GREY)])

members = [
    ("PB","Pushkar Bhangale","Team Leader"),
    ("PA","Prasad Aher","Team Member"),
    ("KT","Kumar Tambe","Team Member"),
    ("AS","Aniruddha Silimkar","Team Member"),
]
cols = [Inches(0.55), Inches(3.85), Inches(7.15), Inches(10.45)]
for i, (init, name, role) in enumerate(members):
    ml = cols[i]; mt = Inches(3.28); mw = Inches(2.9); mh = Inches(0.78)
    card_shape(s, ml, mt, mw, mh)
    av = add_oval(s, ml+Inches(0.12), mt+Inches(0.14), Inches(0.5), Inches(0.5), LBLUE)
    tf2 = av.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = init
    r2.font.name = FONT; r2.font.size = Pt(10); r2.font.bold = True; r2.font.color.rgb = DBLUE
    textbox(s, ml+Inches(0.72), mt+Inches(0.1), mw-Inches(0.82), Inches(0.32),
            [ln(name, 12, bold=True, color=BLACK)])
    textbox(s, ml+Inches(0.72), mt+Inches(0.42), mw-Inches(0.82), Inches(0.26),
            [ln(role, 10, color=GREY)])


# ── SLIDE 2 — Problem ────────────────────────────────────────────
s = blank(prs)
set_bg(s, WHITE)
top_bar(s, RED)
slide_tag(s, "The Problem", LRED, DRED)
slide_title(s, "Cyber attacks are smarter than ever")
slide_num(s, 2, 8)

cw = Inches(3.95); ch = Inches(2.35); ct = Inches(1.62)
stats_data = [
    ("90%",  "Attacks start with a URL",
     "Malicious links are the #1 entry point for hackers",
     "[1] Proofpoint State of the Phish Report, 2023"),
    ("135%", "Rise in AI-based attacks",
     "Increase recorded from 2022 to 2023",
     "[2] Darktrace Annual Threat Report, 2023"),
    ("#1",   "Prompt injection risk",
     "OWASP's top vulnerability for all LLM-based apps",
     "[3] OWASP Top 10 for LLMs — LLM01:2023, owasp.org"),
]
for i, (st, hd, bd, ci) in enumerate(stats_data):
    stat_card(s, Inches(0.55)+i*(cw+Inches(0.19)), ct, cw, ch, st, hd, bd, ci)

divider(s, Inches(4.14))
textbox(s, Inches(0.55), Inches(4.25), Inches(12.2), Inches(0.5),
        [ln("Current tools flag threats but never explain why. Users are left confused. Third Eye solves this.",
            13, color=GREY)])


# ── SLIDE 3 — Solution ───────────────────────────────────────────
s = blank(prs)
set_bg(s, WHITE)
top_bar(s, TEAL)
slide_tag(s, "Our Solution", LTEAL, DTEAL)
slide_title(s, "What is Third Eye?")
slide_num(s, 3, 8)

textbox(s, Inches(0.55), Inches(1.42), Inches(12.2), Inches(0.4),
        [ln("A 3-in-1 cyber threat detection platform that detects, explains, and recommends — in plain English.",
            13, color=GREY)])

cw = Inches(3.95); ch = Inches(2.65); ct = Inches(1.95)
mods = [
    ("Module 1 — URL Scanner",
     ["Paste any link — know instantly if it's phishing,", "malware, defacement, or safe."],
     "[4] Dataset: Kaggle — sid321axn · 641,119 real URLs"),
    ("Module 2 — Prompt Injection Detector",
     ["Catches AI jailbreak & manipulation attempts", "before they cause harm. 4 attack categories."],
     "[5] Dataset: Kaggle — cyberprince Prompt Injection Dataset"),
    ("Module 3 — Risk Dashboard",
     ["Live risk score (0–100), attack category,", "confidence %, and recommended next action."],
     "IndiaNext Problem Statement — Required deliverable"),
]
for i, (hd, bd, ci) in enumerate(mods):
    card_content(s, Inches(0.55)+i*(cw+Inches(0.19)), ct, cw, ch, hd, bd, ci)


# ── SLIDE 4 — Architecture ───────────────────────────────────────
s = blank(prs)
set_bg(s, WHITE)
top_bar(s, BLUE)
slide_tag(s, "Technical Complexity", LBLUE, DBLUE)
slide_title(s, "How it works — system architecture")
slide_num(s, 4, 8)

steps = [
    ("User Input",     "URL or AI Prompt"),
    ("Input Router",   "FastAPI backend"),
    ("Detection Model","RF or DistilBERT"),
    ("Explainability", "Why it's a threat"),
    ("Dashboard",      "Risk score + Action"),
]
bw = Inches(2.12); bh = Inches(0.88); bt = Inches(1.62)
for i, (title, sub) in enumerate(steps):
    bl = Inches(0.55) + i*(bw+Inches(0.22))
    card_shape(s, bl, bt, bw, bh)
    textbox(s, bl+Inches(0.1), bt+Inches(0.08), bw-Inches(0.2), Inches(0.36),
            [ln(title, 12, bold=True, color=BLACK, align=PP_ALIGN.CENTER)])
    textbox(s, bl+Inches(0.1), bt+Inches(0.46), bw-Inches(0.2), Inches(0.3),
            [ln(sub, 10, color=GREY, align=PP_ALIGN.CENTER)])
    if i < 4:
        textbox(s, bl+bw+Inches(0.04), bt+Inches(0.3), Inches(0.18), Inches(0.3),
                [ln("→", 14, color=GREY, align=PP_ALIGN.CENTER)])

divider(s, Inches(2.68))
cw = Inches(3.95); ch = Inches(2.1); ct = Inches(2.82)
arch_cards = [
    ("URL module",
     ["Random Forest · 12 hand-engineered features", "641K URLs · 95–98% accuracy"],
     "[4] Kaggle — sid321axn Malicious URLs Dataset"),
    ("Prompt module",
     ["Fine-tuned DistilBERT · 4 attack categories", "semantic intent understanding"],
     "[5] Kaggle — cyberprince · [7] Sanh et al., arXiv:1910.01108"),
    ("Backend",
     ["FastAPI · separate routes per threat type", "unified response format"],
     "fastapi.tiangolo.com"),
]
for i, (hd, bd, ci) in enumerate(arch_cards):
    card_content(s, Inches(0.55)+i*(cw+Inches(0.19)), ct, cw, ch, hd, bd, ci)


# ── SLIDE 5 — AI/ML ──────────────────────────────────────────────
s = blank(prs)
set_bg(s, WHITE)
top_bar(s, BLUE)
slide_tag(s, "AI / ML Effectiveness", LBLUE, DBLUE)
slide_title(s, "Why we chose these models")
slide_num(s, 5, 8)

hw = Inches(6.05); hh = Inches(4.35); ht = Inches(1.58)
half_card(s, Inches(0.55), ht, hw, hh,
    "URLs  →  Random Forest",
    ["12 hand-engineered URL features (IP, TLD, keywords, length…)",
     "95–98% accuracy on 641,119 labelled URLs",
     "Built-in feature importances = free explainability",
     "sklearn RandomForestClassifier — scikit-learn.org"],
    "[4] Kaggle sid321axn · [6] Pedregosa et al., JMLR 12, 2011")

half_card(s, Inches(6.83), ht, hw, hh,
    "AI Prompts  →  DistilBERT",
    ["distilbert-base-uncased — HuggingFace Transformers",
     "Fine-tuned on real red-team prompt injection dataset",
     "Understands semantic intent, not just keywords",
     "Catches evasion attacks that keyword filters miss"],
    "[5] Kaggle cyberprince · [7] Sanh et al., DistilBERT, arXiv:1910.01108")


# ── SLIDE 6 — Explainability ─────────────────────────────────────
s = blank(prs)
set_bg(s, WHITE)
top_bar(s, PURPLE)
slide_tag(s, "Explainable AI — 15 marks", LPURPLE, DPURPLE)
slide_title(s, "We don't just say 'dangerous' — we explain why")
slide_num(s, 6, 8)

hw = Inches(6.05); hh = Inches(3.72); ht = Inches(1.58)
half_card(s, Inches(0.55), ht, hw, hh,
    "For a malicious URL",
    ["Risk score: 87 / 100  (HIGH THREAT)",
     "Contains keywords: secure, verify",
     "Uses .xyz TLD — common phishing domain",
     "URL is unusually long — obfuscation attempt"],
    "[6] sklearn feature_importances_ — scikit-learn.org")

half_card(s, Inches(6.83), ht, hw, hh,
    "For a prompt injection",
    ["Jailbreak detected — 94% confidence",
     "'Ignore all previous instructions' found",
     "Attempts to override AI safety guidelines",
     "Recommended: Block & log this input"],
    "[5] Category + description fields — cyberprince dataset")

divider(s, Inches(5.46))
textbox(s, Inches(0.55), Inches(5.56), Inches(12.2), Inches(0.45),
        [ln("Every result shows:  what was detected  ·  why it's suspicious  ·  confidence score  ·  what to do next",
            12, color=GREY)])


# ── SLIDE 7 — Innovation ─────────────────────────────────────────
s = blank(prs)
set_bg(s, WHITE)
top_bar(s, AMBER)
slide_tag(s, "Innovation & Trend Alignment", LAMBER, DAMBER)
slide_title(s, "Why Third Eye is future-ready")
slide_num(s, 7, 8)

cw = Inches(5.95); ch = Inches(1.88); gap = Inches(0.2)
inno = [
    ("Multi-threat fusion",
     ["URL + Prompt injection in one platform",
      "Same dashboard, same risk format — not two separate tools"],
     "IndiaNext Problem Statement — 'Hybrid cyber threat intelligence dashboard'"),
    ("Defending AI with AI",
     ["Transformer model protects AI systems from being attacked",
      "The frontier of cybersecurity in 2025"],
     "[3] OWASP LLM01:2023 · [9] IBM X-Force Threat Intelligence Index, 2024"),
    ("Adversarial testing",
     ["We red-teamed our own classifier with evasion prompts",
      "It still catches them — robustness proven"],
     "IndiaNext Bonus Criteria — 'Adversarial robustness testing'"),
    ("Market opportunity",
     ["AI security tools market projected at $102B by 2032",
      "We are building the defence now"],
     "[10] MarketsandMarkets AI Security Report, 2023"),
]
for i, (hd, bd, ci) in enumerate(inno):
    col_l = Inches(0.55) if i % 2 == 0 else Inches(0.55)+cw+gap
    row_t = Inches(1.58) if i < 2 else Inches(1.58)+ch+gap
    card_content(s, col_l, row_t, cw, ch, hd, bd, ci)


# ── SLIDE 8 — Closing ────────────────────────────────────────────
s = blank(prs)
set_bg(s, WHITE)
top_bar(s, TEAL)
slide_tag(s, "Closing", LTEAL, DTEAL)
slide_title(s, "Third Eye — built for the AI era of cyber threats")
slide_num(s, 8, 8)

textbox(s, Inches(0.55), Inches(1.42), Inches(12.2), Inches(0.4),
        [ln("Not defending against 2010 threats. Defending against 2025 threats — attacks on AI itself.",
            13, color=GREY)])

cw = Inches(3.95); ch = Inches(1.52); ct = Inches(1.95)
closing_stats = [
    ("641K",   "Real URLs trained on",    "[4] Kaggle sid321axn"),
    ("95–98%", "URL detection accuracy",  "[4] sklearn RF classifier"),
    ("$102B",  "AI security market 2032", "[10] MarketsandMarkets, 2023"),
]
for i, (st, bd, ci) in enumerate(closing_stats):
    stat_card(s, Inches(0.55)+i*(cw+Inches(0.19)), ct, cw, ch, st, bd, "", ci)

divider(s, Inches(3.64))

refs = (
    "[1] Proofpoint State of the Phish, 2023  ·  [2] Darktrace Annual Threat Report, 2023  ·  "
    "[3] OWASP Top 10 for LLMs, LLM01:2023  ·  [4] Kaggle: sid321axn/malicious-urls-dataset  ·  "
    "[5] Kaggle: cyberprince/prompt-injection-and-benign-prompt-dataset  ·  "
    "[6] Scikit-learn, JMLR 12, 2011  ·  [7] Sanh et al., DistilBERT, arXiv:1910.01108  ·  "
    "[8] McKinsey State of AI, 2024  ·  [9] IBM X-Force Index, 2024  ·  "
    "[10] MarketsandMarkets AI Security, 2023"
)
card_shape(s, Inches(0.55), Inches(3.76), Inches(12.2), Inches(0.82))
textbox(s, Inches(0.72), Inches(3.84), Inches(11.88), Inches(0.68),
        [ln(refs, 8, italic=True, color=GREY)])

badges = ["OWASP LLM01:2023","Random Forest","DistilBERT","FastAPI","Explainable AI","Multi-threat fusion"]
bx = Inches(0.55)
for badge in badges:
    bw2 = Inches(1.75); bh2 = Inches(0.3)
    shp = add_rect(s, bx, Inches(4.72), bw2, bh2, fill_color=LPURPLE)
    shp.line.fill.background()
    tf = shp.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = badge
    r.font.name = FONT; r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = DPURPLE
    bx += bw2 + Inches(0.12)

textbox(s, Inches(0.55), Inches(5.18), Inches(12.2), Inches(0.42),
        [ln("Team Runtime Errors  ·  K.E.S. Shroff College  ·  IndiaNext Hackathon 2026",
            13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)])

# ── Save ──────────────────────────────────────────────────────────
prs.save("Third_Eye_RuntimeErrors.pptx")
print("✅  Saved → Third_Eye_RuntimeErrors.pptx")